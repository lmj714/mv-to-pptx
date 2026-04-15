import streamlit as st
import yt_dlp
import os
import re
import time
import tempfile
import io
import cv2
import requests
import zipfile
import shutil
import numpy as np
from pathlib import Path
from PIL import Image, ImageEnhance, ImageFilter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from deep_translator import GoogleTranslator
from gtts import gTTS

# ── Page config ─────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="K-POP 韓語學習簡報產生器",
    page_icon="🎤",
    layout="centered",
)

# ── Custom CSS ───────────────────────────────────────────────────────────────
st.markdown("""
<style>
    .main-title {
        font-size: 2.2rem;
        font-weight: 800;
        text-align: center;
        background: linear-gradient(135deg, #FF6B9D, #C44B8A, #9B59B6);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        margin-bottom: 0.2rem;
    }
    .sub-title {
        text-align: center;
        color: #888;
        font-size: 0.95rem;
        margin-bottom: 2rem;
    }
    .stButton > button {
        width: 100%;
        background: linear-gradient(135deg, #FF6B9D, #9B59B6);
        color: white;
        border: none;
        border-radius: 8px;
        padding: 0.6rem 1.2rem;
        font-size: 1rem;
        font-weight: 600;
        cursor: pointer;
    }
    .stButton > button:hover { opacity: 0.9; }
</style>
""", unsafe_allow_html=True)

# ── Title ────────────────────────────────────────────────────────────────────
st.markdown('<div class="main-title">K-POP 韓語學習簡報產生器 🎤</div>', unsafe_allow_html=True)
st.markdown('<div class="sub-title">貼上 YouTube MV 網址，自動產生韓中對照投影片（含 MV 畫面）</div>', unsafe_allow_html=True)

# ── Input ────────────────────────────────────────────────────────────────────
url = st.text_input(
    "YouTube 網址",
    placeholder="https://www.youtube.com/watch?v=...",
    label_visibility="collapsed",
)

import socket
def _is_cloud() -> bool:
    """Detect if running on Streamlit Cloud (AWS) by checking hostname."""
    try:
        h = socket.gethostname()
        return "streamlit" in h.lower() or h.startswith("ip-")
    except Exception:
        return False

_running_on_cloud = _is_cloud()

col1, col2 = st.columns([1, 2])
with col1:
    if _running_on_cloud:
        use_mv_frames = False
        st.info("☁️ 雲端版使用縮圖背景（YouTube 限制雲端下載）")
    else:
        use_mv_frames = st.checkbox("🎬 擷取 MV 畫面", value=True,
                                    help="下載低畫質影片並擷取對應畫面作為背景（較慢）；取消勾選則改用縮圖，速度快很多")
with col2:
    use_tts = st.checkbox("🔊 加入韓文語音（TTS）", value=True,
                          help="為每張投影片加入 AI 朗讀的韓文語音，在 PowerPoint 中點擊即可播放")

# ═══════════════════════════════════════════════════════════════════════════════
# Helper functions
# ═══════════════════════════════════════════════════════════════════════════════

def get_video_id(url: str) -> str | None:
    m = re.search(r"(?:v=|youtu\.be/)([a-zA-Z0-9_-]{11})", url)
    return m.group(1) if m else None


def get_video_info(url: str) -> dict:
    ydl_opts = {
        "quiet": True,
        "no_warnings": True,
        "skip_download": True,
        "socket_timeout": 15,
    }
    try:
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(url, download=False)
            if info is None:
                raise ValueError("影片不存在或無法存取")
            return info
    except SystemExit as e:
        raise ValueError(f"yt-dlp 錯誤（影片可能不存在或有地區限制）：{e}")


def download_subtitles(url: str, tmpdir: str) -> tuple[str | None, str]:
    """Try manual then auto Korean subtitles. Returns (path, info_str)."""
    base = os.path.join(tmpdir, "subtitle")
    for auto in (False, True):
        ydl_opts = {
            "quiet": True, "no_warnings": True, "skip_download": True,
            "writesubtitles": not auto, "writeautomaticsub": auto,
            "subtitleslangs": ["ko"], "subtitlesformat": "vtt",
            "outtmpl": base,
        }
        try:
            with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                ydl.download([url])
        except (SystemExit, Exception):
            pass
        for ext in ("vtt", "srt"):
            candidate = f"{base}.ko.{ext}"
            if os.path.exists(candidate):
                return candidate, "自動生成" if auto else "官方提供"
    return None, "找不到韓文字幕"


def _ts_to_sec(ts: str) -> float:
    """Convert HH:MM:SS.mmm or MM:SS.mmm to seconds."""
    parts = ts.strip().split(":")
    if len(parts) == 3:
        h, m, s = int(parts[0]), int(parts[1]), float(parts[2])
    else:
        h, m, s = 0, int(parts[0]), float(parts[1])
    return h * 3600 + m * 60 + s


def parse_vtt(path: str) -> list[tuple[float, str]]:
    """Return [(start_sec, korean_text), ...] from .vtt file."""
    with open(path, encoding="utf-8") as f:
        content = f.read()

    results: list[tuple[float, str]] = []
    last_text: str = ""   # Only skip CONSECUTIVE duplicates (subtitle overlap)

    for block in re.split(r"\n\s*\n", content.strip()):
        lines = block.strip().splitlines()
        ts_sec = None
        text_parts: list[str] = []

        for line in lines:
            ts_m = re.match(r"([\d:]+\.[\d]+)\s*-->", line)
            if ts_m:
                ts_sec = _ts_to_sec(ts_m.group(1))
                continue
            if re.match(r"^\d+$", line.strip()) or re.match(r"WEBVTT|NOTE", line):
                continue
            clean = re.sub(r"<[^>]+>", "", line).strip()
            if clean:
                text_parts.append(clean)

        if ts_sec is not None and text_parts:
            text = " ".join(text_parts)
            if text != last_text:   # Allow repeats (chorus), block only immediate dupes
                last_text = text
                results.append((ts_sec, text))

    return results


def parse_srt(path: str) -> list[tuple[float, str]]:
    """Return [(start_sec, korean_text), ...] from .srt file."""
    with open(path, encoding="utf-8") as f:
        content = f.read()

    results: list[tuple[float, str]] = []
    last_text: str = ""   # Only skip CONSECUTIVE duplicates

    for block in re.split(r"\n\s*\n", content.strip()):
        lines = block.strip().splitlines()
        ts_sec = None
        text_parts: list[str] = []

        for line in lines:
            if re.match(r"^\d+$", line.strip()):
                continue
            ts_m = re.match(r"([\d:]+,[\d]+)\s*-->", line)
            if ts_m:
                ts_sec = _ts_to_sec(ts_m.group(1).replace(",", "."))
                continue
            clean = re.sub(r"<[^>]+>", "", line).strip()
            if clean:
                text_parts.append(clean)

        if ts_sec is not None and text_parts:
            text = " ".join(text_parts)
            if text != last_text:   # Allow repeats (chorus), block only immediate dupes
                last_text = text
                results.append((ts_sec, text))

    return results


def translate_with_retry(text: str, retries: int = 4, delay: float = 1.5) -> str:
    """Translate Korean → Traditional Chinese with retry on failure."""
    for attempt in range(retries):
        try:
            result = GoogleTranslator(source="ko", target="zh-TW").translate(text)
            if result:
                return result
        except Exception:
            pass
        time.sleep(delay * (attempt + 1))
    return ""   # Return empty string only after all retries exhausted


def translate_lines(entries: list[tuple[float, str]],
                    status_ph) -> list[str]:
    """Translate all Korean lines, showing progress."""
    translated: list[str] = []
    total = len(entries)

    for i, (_, ko) in enumerate(entries, 1):
        status_ph.text(f"翻譯中... ({i}/{total}）")
        zh = translate_with_retry(ko)
        translated.append(zh)
        time.sleep(0.1)   # Minimal delay to avoid rate limiting

    return translated


def _ffmpeg_available() -> bool:
    """Check if ffmpeg is on PATH."""
    import shutil
    return shutil.which("ffmpeg") is not None


def download_video(url: str, tmpdir: str, status_ph) -> str | None:
    """Download lowest-quality video for frame extraction. Returns file path."""
    out_tmpl = os.path.join(tmpdir, "video.%(ext)s")
    fmt = (
        "bestvideo[height<=720][ext=mp4]"
        "/bestvideo[height<=720]"
        "/bestvideo[height<=480][ext=mp4]"
        "/bestvideo[height<=480]"
        "/worst[ext=mp4]/worst"
    )

    ydl_opts = {
        "quiet": True,
        "no_warnings": True,
        "format": fmt,
        "outtmpl": out_tmpl,
        "noplaylist": True,
    }
    status_ph.text("正在下載 MV 影片（低畫質）以擷取畫面…")
    try:
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            ydl.extract_info(url, download=True)
        for f in os.listdir(tmpdir):
            if f.startswith("video."):
                return os.path.join(tmpdir, f)
    except Exception as e:
        status_ph.text(f"影片下載失敗（將改用縮圖背景）：{e}")
    return None


def get_thumbnail(video_id: str) -> bytes | None:
    """Download YouTube thumbnail, try maxres then hq."""
    for quality in ("maxresdefault", "hqdefault", "mqdefault"):
        try:
            resp = requests.get(
                f"https://img.youtube.com/vi/{video_id}/{quality}.jpg",
                timeout=10
            )
            if resp.status_code == 200 and len(resp.content) > 5000:
                return resp.content
        except Exception:
            pass
    return None


def _process_frame_image(img: Image.Image, brightness: float = 0.60) -> bytes:
    """Resize to 1280×720 (PPT 13.33"×7.5" @ 96dpi), brighten, light blur → JPEG."""
    img = img.convert("RGB").resize((1280, 720), Image.LANCZOS)
    img = ImageEnhance.Brightness(img).enhance(brightness)
    img = img.filter(ImageFilter.GaussianBlur(radius=0.8))
    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=92)
    buf.seek(0)
    return buf.read()


def extract_all_frames(video_path: str, timestamps: list[float],
                       brightness: float = 0.60,
                       status_ph=None) -> dict[float, bytes]:
    """
    Open the video ONCE with cv2 and extract all frames in a single pass.
    Returns {timestamp: jpeg_bytes}.  Much faster than per-frame subprocess.
    """
    frames: dict[float, bytes] = {}
    if not video_path or not os.path.exists(video_path):
        return frames

    try:
        cap = cv2.VideoCapture(video_path)
        total = len(timestamps)
        for i, ts in enumerate(timestamps):
            if status_ph:
                status_ph.text(f"擷取 MV 畫面... ({i+1}/{total})")
            cap.set(cv2.CAP_PROP_POS_MSEC, ts * 1000)
            ret, frame_bgr = cap.read()
            if ret:
                frame_rgb = cv2.cvtColor(frame_bgr, cv2.COLOR_BGR2RGB)
                frames[ts] = _process_frame_image(Image.fromarray(frame_rgb), brightness)
        cap.release()
    except Exception:
        pass

    return frames


def generate_tts_audio(text: str, lang: str = "ko") -> bytes | None:
    """Generate TTS mp3 bytes for given text using gTTS."""
    try:
        buf = io.BytesIO()
        gTTS(text=text, lang=lang, slow=False).write_to_fp(buf)
        buf.seek(0)
        return buf.read()
    except Exception:
        return None


def generate_all_tts(entries: list[tuple[float, str]], status_ph) -> dict[int, bytes]:
    """Generate TTS for all Korean lyrics. Returns {index: mp3_bytes}."""
    result: dict[int, bytes] = {}
    total = len(entries)
    for i, (_, ko) in enumerate(entries):
        status_ph.text(f"生成語音... ({i+1}/{total})")
        audio = generate_tts_audio(ko, lang="ko")
        if audio:
            result[i] = audio
    return result


def embed_audio_into_pptx(pptx_bytes: bytes,
                           audio_map: dict[int, bytes]) -> bytes:
    """
    Embed TTS mp3 files into each lyric slide (slide index 1..N, slide 0 is cover).
    Uses direct ZIP manipulation since python-pptx has no audio API.
    Returns modified pptx bytes.
    """
    if not audio_map:
        return pptx_bytes

    with tempfile.TemporaryDirectory() as tmp:
        # ── Unpack PPTX (it's a ZIP) ─────────────────────────────────────────
        pptx_path = os.path.join(tmp, "deck.pptx")
        with open(pptx_path, "wb") as f:
            f.write(pptx_bytes)

        unpack_dir = os.path.join(tmp, "unpacked")
        with zipfile.ZipFile(pptx_path, "r") as z:
            z.extractall(unpack_dir)

        media_dir = os.path.join(unpack_dir, "ppt", "media")
        os.makedirs(media_dir, exist_ok=True)
        slides_dir  = os.path.join(unpack_dir, "ppt", "slides")
        rels_dir    = os.path.join(unpack_dir, "ppt", "slides", "_rels")
        os.makedirs(rels_dir, exist_ok=True)

        # ── Embed audio for each lyric slide ────────────────────────────────
        # Slide 0 = cover → lyric slides start at slide index 1
        for lyric_idx, mp3_bytes in audio_map.items():
            slide_num = lyric_idx + 2          # slide files are 1-based; +1 for cover
            media_name = f"audio{lyric_idx + 1}.mp3"
            media_path = os.path.join(media_dir, media_name)

            # Write mp3
            with open(media_path, "wb") as f:
                f.write(mp3_bytes)

            # ── Update slide rels XML ────────────────────────────────────────
            rels_file = os.path.join(rels_dir, f"slide{slide_num}.xml.rels")
            audio_rel_type = (
                "http://schemas.openxmlformats.org/officeDocument/2006/"
                "relationships/audio"
            )

            if os.path.exists(rels_file):
                with open(rels_file, encoding="utf-8") as f:
                    rels_xml = f.read()
                # Find next rId number
                existing_ids = re.findall(r'Id="rId(\d+)"', rels_xml)
                next_id = max((int(x) for x in existing_ids), default=0) + 1
                audio_rid = f"rId{next_id}"
                # Insert before closing tag
                new_rel = (
                    f'<Relationship Id="{audio_rid}" '
                    f'Type="{audio_rel_type}" '
                    f'Target="../media/{media_name}"/>'
                )
                rels_xml = rels_xml.replace("</Relationships>",
                                            f"{new_rel}</Relationships>")
            else:
                audio_rid = "rId1"
                rels_xml = (
                    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
                    '<Relationships xmlns="http://schemas.openxmlformats.org/'
                    'package/2006/relationships">'
                    f'<Relationship Id="{audio_rid}" '
                    f'Type="{audio_rel_type}" '
                    f'Target="../media/{media_name}"/>'
                    "</Relationships>"
                )

            with open(rels_file, "w", encoding="utf-8") as f:
                f.write(rels_xml)

            # ── Update slide XML: add audio element ──────────────────────────
            slide_file = os.path.join(slides_dir, f"slide{slide_num}.xml")
            if not os.path.exists(slide_file):
                continue

            with open(slide_file, encoding="utf-8") as f:
                slide_xml = f.read()

            # Build audio shape XML (invisible icon, auto-play, loop-off)
            audio_shape = f"""<p:sp>
  <p:nvSpPr>
    <p:cNvPr id="999" name="Audio_{lyric_idx}"/>
    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
    <p:nvPr>
      <p:ph type="body"/>
      <a:audioFile r:link="{audio_rid}"/>
    </p:nvPr>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="457200" y="6400000"/><a:ext cx="457200" cy="457200"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
</p:sp>"""

            # PowerPoint requires a <p:timing> block to auto-play audio
            timing_xml = f"""<p:timing>
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefin" restart="whenNotActive" nodeType="tmRoot">
        <p:childTnLst>
          <p:par>
            <p:cTn id="2" fill="hold">
              <p:stCondLst><p:cond delay="0"/></p:stCondLst>
              <p:childTnLst>
                <p:par>
                  <p:cTn id="3" presetID="1" presetClass="mediacall"
                         presetSubtype="0" fill="hold" nodeType="clickEffect">
                    <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                    <p:childTnLst>
                      <p:audio>
                        <p:cMediaNode vol="80000" showWhenStopped="0">
                          <p:cTn id="4" fill="hold"><p:stCondLst>
                            <p:cond delay="0"/>
                          </p:stCondLst></p:cTn>
                          <p:tgtEl>
                            <p:spTgt spid="999"/>
                          </p:tgtEl>
                        </p:cMediaNode>
                      </p:audio>
                    </p:childTnLst>
                  </p:cTn>
                </p:par>
              </p:childTnLst>
            </p:cTn>
          </p:par>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
</p:timing>"""

            # Insert audio shape before </p:spTree>
            slide_xml = slide_xml.replace("</p:spTree>",
                                          f"{audio_shape}</p:spTree>")
            # Insert timing before </p:cSld> closing
            if "<p:timing>" not in slide_xml:
                slide_xml = slide_xml.replace("</p:cSld>",
                                              f"</p:cSld>{timing_xml}")

            with open(slide_file, "w", encoding="utf-8") as f:
                f.write(slide_xml)

        # ── Re-pack into PPTX ────────────────────────────────────────────────
        out_path = os.path.join(tmp, "output.pptx")
        with zipfile.ZipFile(out_path, "w", zipfile.ZIP_DEFLATED) as zout:
            for root, _, files in os.walk(unpack_dir):
                for fname in files:
                    fpath = os.path.join(root, fname)
                    arcname = os.path.relpath(fpath, unpack_dir)
                    zout.write(fpath, arcname)

        with open(out_path, "rb") as f:
            return f.read()


def prepare_thumbnail_bg(thumb_bytes: bytes | None,
                         brightness: float = 0.60) -> bytes | None:
    """Darken thumbnail for use as slide background."""
    if not thumb_bytes:
        return None
    try:
        img = Image.open(io.BytesIO(thumb_bytes)).convert("RGB")
        img = img.resize((1280, 720), Image.LANCZOS)
        img = ImageEnhance.Brightness(img).enhance(brightness)
        img = img.filter(ImageFilter.GaussianBlur(radius=1.5))
        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=82)
        buf.seek(0)
        return buf.read()
    except Exception:
        return None


def hex_color(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _add_bg_image(slide, img_bytes: bytes, W, H):
    """Add image as full-slide background (z-order: bottom)."""
    pic = slide.shapes.add_picture(io.BytesIO(img_bytes), 0, 0, W, H)
    sp_tree = slide.shapes._spTree
    sp_tree.remove(pic._element)
    sp_tree.insert(2, pic._element)   # Push to back


def _add_solid_bg(slide, color: RGBColor, W, H):
    """Fallback: solid color background rectangle."""
    shape = slide.shapes.add_shape(1, 0, 0, W, H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_text(slide, text, left, top, width, height,
              size, bold, color, align=PP_ALIGN.CENTER, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def build_pptx(
    title: str,
    entries: list[tuple[float, str]],
    chinese_lines: list[str],
    thumb_bytes: bytes | None,
    frames: dict[float, bytes],   # pre-extracted: {timestamp: jpeg_bytes}
    status_ph,
) -> bytes:
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    W, H = prs.slide_width, prs.slide_height

    COLOR_BG      = hex_color("1A0533")
    COLOR_ACCENT  = hex_color("FF6B9D")
    COLOR_WHITE   = hex_color("FFFFFF")
    COLOR_CHINESE = hex_color("F9C6DD")
    COLOR_SUB     = hex_color("C490D1")

    blank = prs.slide_layouts[6]

    # ── Cover slide ──────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank)

    # Background: thumbnail (full brightness for cover) or solid
    if thumb_bytes:
        try:
            img = Image.open(io.BytesIO(thumb_bytes)).convert("RGB")
            img = img.resize((1280, 720), Image.LANCZOS)
            img = ImageEnhance.Brightness(img).enhance(0.50)
            img = img.filter(ImageFilter.GaussianBlur(radius=2))
            buf = io.BytesIO()
            img.save(buf, format="JPEG", quality=82)
            buf.seek(0)
            _add_bg_image(slide, buf.read(), W, H)
        except Exception:
            _add_solid_bg(slide, COLOR_BG, W, H)
    else:
        _add_solid_bg(slide, COLOR_BG, W, H)

    # Top bar
    bar = slide.shapes.add_shape(1, 0, 0, W, Inches(0.07))
    bar.fill.solid(); bar.fill.fore_color.rgb = COLOR_ACCENT; bar.line.fill.background()

    _add_text(slide, "K-POP 韓語學習",
              Inches(0.5), Inches(1.8), Inches(12.33), Inches(0.6),
              16, False, COLOR_SUB, italic=True)
    _add_text(slide, title,
              Inches(0.5), Inches(2.5), Inches(12.33), Inches(2.5),
              40, True, COLOR_WHITE)
    _add_text(slide, "· · ·",
              Inches(0.5), Inches(5.1), Inches(12.33), Inches(0.5),
              20, False, COLOR_ACCENT)
    _add_text(slide, "韓中對照學習投影片  /  한중 대조 학습 슬라이드",
              Inches(0.5), Inches(5.7), Inches(12.33), Inches(0.6),
              14, False, COLOR_SUB)

    bot = slide.shapes.add_shape(1, 0, H - Inches(0.07), W, Inches(0.07))
    bot.fill.solid(); bot.fill.fore_color.rgb = COLOR_ACCENT; bot.line.fill.background()

    # ── Thumbnail fallback background (darkened) ─────────────────────────────
    thumb_bg = prepare_thumbnail_bg(thumb_bytes)

    # ── Lyric slides ─────────────────────────────────────────────────────────
    total = len(entries)
    for idx, ((ts, ko), zh) in enumerate(zip(entries, chinese_lines), 1):
        status_ph.text(f"正在生成投影片... ({idx}/{total})")
        slide = prs.slides.add_slide(blank)

        # Use pre-extracted frame; fallback to thumbnail
        frame_bg = frames.get(ts)
        if frame_bg:
            _add_bg_image(slide, frame_bg, W, H)
        elif thumb_bg:
            _add_bg_image(slide, thumb_bg, W, H)
        else:
            _add_solid_bg(slide, COLOR_BG, W, H)

        # Top accent bar
        t = slide.shapes.add_shape(1, 0, 0, W, Inches(0.06))
        t.fill.solid(); t.fill.fore_color.rgb = COLOR_ACCENT; t.line.fill.background()

        # Slide number
        _add_text(slide, f"{idx:02d}",
                  Inches(0.4), Inches(0.15), Inches(1), Inches(0.5),
                  13, True, COLOR_ACCENT, align=PP_ALIGN.LEFT)

        # Korean lyric
        _add_text(slide, ko,
                  Inches(0.8), Inches(1.8), Inches(11.73), Inches(2.4),
                  34, True, COLOR_WHITE)

        # Separator dot
        sep = slide.shapes.add_shape(1, Inches(4.5), Inches(4.35), Inches(4.33), Pt(1))
        sep.fill.solid(); sep.fill.fore_color.rgb = COLOR_ACCENT; sep.line.fill.background()

        # Chinese translation
        _add_text(slide, zh if zh else ko,   # If translation failed, show Korean again
                  Inches(0.8), Inches(4.6), Inches(11.73), Inches(1.8),
                  22, False, COLOR_CHINESE)

        # Bottom bar
        b = slide.shapes.add_shape(1, 0, H - Inches(0.06), W, Inches(0.06))
        b.fill.solid(); b.fill.fore_color.rgb = COLOR_ACCENT; b.line.fill.background()

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ═══════════════════════════════════════════════════════════════════════════════
# Main flow
# ═══════════════════════════════════════════════════════════════════════════════

if st.button("🎵 開始轉換"):
    if not url.strip():
        st.error("請輸入 YouTube 網址。"); st.stop()
    if not re.match(r"https?://(www\.)?(youtube\.com|youtu\.be)/", url.strip()):
        st.error("請輸入有效的 YouTube 網址。"); st.stop()

    # 移除播放清單參數，只保留 v=VIDEO_ID
    vid = get_video_id(url.strip())
    if vid:
        url = f"https://www.youtube.com/watch?v={vid}"

    pptx_bytes: bytes | None = None
    video_title = "K-POP 影片"

    with tempfile.TemporaryDirectory() as tmpdir:

        # ── Step 1: Video info ───────────────────────────────────────────────
        with st.spinner("正在取得影片資訊…"):
            try:
                info = get_video_info(url.strip())
                video_title = info.get("title", "K-POP 影片")
                video_id    = get_video_id(url.strip()) or info.get("id", "")
                st.info(f"影片標題：**{video_title}**")
            except Exception as e:
                st.error(f"無法取得影片資訊：{e}"); st.stop()

        # ── Step 2: Thumbnail ────────────────────────────────────────────────
        thumb_bytes: bytes | None = None
        with st.spinner("正在下載縮圖…"):
            thumb_bytes = get_thumbnail(video_id)
            if thumb_bytes:
                st.success("縮圖下載完成")

        # ── Step 3: Subtitles ────────────────────────────────────────────────
        with st.spinner("正在下載韓文字幕…"):
            sub_path, sub_info = download_subtitles(url.strip(), tmpdir)
            if not sub_path:
                st.error(f"字幕下載失敗：{sub_info}\n此影片可能沒有韓文字幕。"); st.stop()
            st.success(f"字幕下載成功（{sub_info}）")

        # ── Step 4: Parse subtitles ──────────────────────────────────────────
        with st.spinner("正在解析字幕…"):
            try:
                entries = (parse_vtt(sub_path) if sub_path.endswith(".vtt")
                           else parse_srt(sub_path))
                if not entries:
                    st.error("字幕解析後沒有找到文字，請換一部影片。"); st.stop()
                st.success(f"共解析到 {len(entries)} 句歌詞")
            except Exception as e:
                st.error(f"字幕解析失敗：{e}"); st.stop()

        # ── Step 5: Download video (optional) ───────────────────────────────
        video_path: str | None = None
        frames: dict[float, bytes] = {}

        if use_mv_frames:
            dl_status = st.empty()
            with st.spinner("正在下載 MV 影片（144p）…"):
                video_path = download_video(url.strip(), tmpdir, dl_status)
                dl_status.empty()
                if video_path:
                    st.success("影片下載完成")
                else:
                    st.warning("影片下載失敗，改用縮圖背景")
        else:
            st.info("已略過影片下載，使用縮圖作為背景")

        # ── Step 6: Translate + extract frames ──────────────────────────────
        trans_status = st.empty()

        # 6a. Batch frame extraction
        if video_path:
            frame_status = st.empty()
            with st.spinner("正在批次擷取 MV 畫面…"):
                timestamps = [ts for ts, _ in entries]
                frames = extract_all_frames(video_path, timestamps, status_ph=frame_status)
                frame_status.empty()
                st.success(f"擷取 {len(frames)}/{len(entries)} 張 MV 畫面")

        # 6b. Translate
        with st.spinner("正在翻譯成繁體中文…"):
            try:
                chinese_lines = translate_lines(entries, trans_status)
                trans_status.empty()
                failed = sum(1 for z in chinese_lines if not z)
                if failed:
                    st.warning(f"翻譯完成（{failed} 句保留韓文原文）")
                else:
                    st.success("翻譯完成！")
            except Exception as e:
                st.error(f"翻譯失敗：{e}"); st.stop()

        # ── Step 7: Build PPTX ───────────────────────────────────────────────
        pptx_status = st.empty()
        with st.spinner("正在生成投影片…"):
            try:
                pptx_bytes = build_pptx(
                    video_title, entries, chinese_lines,
                    thumb_bytes, frames, pptx_status
                )
                pptx_status.empty()
                st.success(f"投影片生成完成！共 {len(entries) + 1} 頁（含封面）")
            except Exception as e:
                st.error(f"投影片生成失敗：{e}"); st.stop()

        # ── Step 8: TTS（optional）──────────────────────────────────────────
        if use_tts and pptx_bytes:
            tts_status = st.empty()
            with st.spinner("正在生成韓文語音（TTS）…"):
                try:
                    audio_map = generate_all_tts(entries, tts_status)
                    tts_status.empty()
                    if audio_map:
                        pptx_bytes = embed_audio_into_pptx(pptx_bytes, audio_map)
                        st.success(f"語音嵌入完成！共 {len(audio_map)} 句")
                    else:
                        st.warning("TTS 生成失敗，投影片不含語音")
                except Exception as e:
                    tts_status.empty()
                    st.warning(f"TTS 嵌入失敗（投影片仍可下載）：{e}")

    # ── Step 8: Download ─────────────────────────────────────────────────────
    if pptx_bytes:
        safe = re.sub(r'[\\/*?:"<>|]', "_", video_title)[:50]
        st.download_button(
            label="⬇️ 下載投影片 (.pptx)",
            data=pptx_bytes,
            file_name=f"{safe}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

# ── Footer ───────────────────────────────────────────────────────────────────
st.markdown("---")
st.markdown(
    "<p style='text-align:center;color:#888;font-size:0.8rem;'>"
    "僅供個人韓語學習用途 · For personal language learning only"
    "</p>",
    unsafe_allow_html=True,
)
